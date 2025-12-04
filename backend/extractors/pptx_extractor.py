"""
PPTX (PowerPoint) text extraction module.
Extracts slide content, titles, and speaker notes.
"""
import logging
from pathlib import Path
from typing import Dict, List, Tuple
from pptx import Presentation
from pptx.util import Inches

logger = logging.getLogger("pinguin_backend.extractors.pptx")


class PPTXExtractor:
    """
    Extracts text from PowerPoint presentations.
    Each slide is treated as a semantic unit with title, content, and notes.
    """
    
    def __init__(self):
        """Initialize PPTX extractor."""
        self.logger = logger
    
    def extract(self, file_path: str) -> Tuple[str, Dict]:
        """
        Extract text and metadata from a PPTX file.
        
        Args:
            file_path: Path to the PPTX file
        
        Returns:
            Tuple of (extracted_text, metadata)
            
        Raises:
            FileNotFoundError: If file doesn't exist
            ValueError: If file is invalid
        """
        path = Path(file_path)
        
        if not path.exists():
            raise FileNotFoundError(f"PPTX file not found: {file_path}")
        
        if not path.suffix.lower() in ['.pptx', '.ppt']:
            raise ValueError(f"File is not a PowerPoint presentation: {file_path}")
        
        self.logger.info(f"Extracting text from PPTX: {path.name}")
        
        try:
            prs = Presentation(file_path)
        except Exception as e:
            self.logger.error(f"Error opening PPTX file: {e}")
            raise ValueError(f"Failed to open PPTX file: {str(e)}")
        
        # Extract text from all slides
        slide_texts = []
        slides_with_notes = 0
        total_shapes = 0
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_data = self._process_slide(slide, slide_num)
            slide_texts.append(slide_data)
            
            if slide_data['has_notes']:
                slides_with_notes += 1
            
            total_shapes += slide_data['shape_count']
        
        # Combine all slide texts
        full_text_parts = []
        for slide_data in slide_texts:
            slide_text = self._format_slide_text(slide_data)
            if slide_text:
                full_text_parts.append(slide_text)
        
        full_text = "\n\n".join(full_text_parts)
        
        metadata = {
            'filename': path.name,
            'file_type': 'pptx',
            'extractor': 'python-pptx',
            'slide_count': len(prs.slides),
            'slides_with_notes': slides_with_notes,
            'total_shapes': total_shapes,
            'slides': slide_texts,
            'char_count': len(full_text)
        }
        
        self.logger.info(
            f"Extracted {len(full_text)} characters from {len(prs.slides)} slides. "
            f"Slides with notes: {slides_with_notes}"
        )
        
        return full_text, metadata
    
    def _process_slide(self, slide, slide_number: int) -> Dict:
        """
        Process a single slide and extract all text content.
        
        Args:
            slide: Slide object from python-pptx
            slide_number: Slide number (1-indexed)
            
        Returns:
            Dictionary with slide data
        """
        slide_data = {
            'slide_number': slide_number,
            'title': '',
            'content': [],
            'notes': '',
            'has_notes': False,
            'shape_count': 0
        }
        
        # Extract text from shapes
        for shape in slide.shapes:
            slide_data['shape_count'] += 1
            
            # Check if shape has text
            if not hasattr(shape, "text"):
                continue
            
            text = shape.text.strip()
            if not text:
                continue
            
            # Check if it's a title
            if shape.is_placeholder:
                placeholder = shape.placeholder_format
                if placeholder.type == 1:  # Title placeholder
                    slide_data['title'] = text
                    continue
            
            # Otherwise, it's content
            slide_data['content'].append(text)
        
        # Extract speaker notes
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            notes_text = notes_text_frame.text.strip()
            
            if notes_text:
                slide_data['notes'] = notes_text
                slide_data['has_notes'] = True
        
        return slide_data
    
    def _format_slide_text(self, slide_data: Dict) -> str:
        """
        Format slide data into readable text.
        
        Args:
            slide_data: Dictionary with slide information
            
        Returns:
            Formatted slide text
        """
        parts = []
        
        # Add slide header
        parts.append(f"[Slide {slide_data['slide_number']}]")
        
        # Add title
        if slide_data['title']:
            parts.append(f"Title: {slide_data['title']}")
        
        # Add content
        if slide_data['content']:
            parts.append("Content:")
            for content_item in slide_data['content']:
                parts.append(f"  {content_item}")
        
        # Add notes
        if slide_data['has_notes']:
            parts.append(f"Notes: {slide_data['notes']}")
        
        return "\n".join(parts)
    
    def extract_slide(self, file_path: str, slide_number: int) -> str:
        """
        Extract text from a specific slide.
        
        Args:
            file_path: Path to PPTX file
            slide_number: Slide number (1-indexed)
            
        Returns:
            Extracted text from the slide
        """
        try:
            prs = Presentation(file_path)
            
            if slide_number < 1 or slide_number > len(prs.slides):
                raise ValueError(
                    f"Slide number {slide_number} out of range "
                    f"(1-{len(prs.slides)})"
                )
            
            slide = prs.slides[slide_number - 1]
            slide_data = self._process_slide(slide, slide_number)
            
            return self._format_slide_text(slide_data)
        
        except Exception as e:
            self.logger.error(f"Error extracting slide {slide_number}: {e}")
            raise
    
    def get_slide_titles(self, file_path: str) -> List[str]:
        """
        Extract all slide titles from presentation.
        
        Args:
            file_path: Path to PPTX file
            
        Returns:
            List of slide titles
        """
        try:
            prs = Presentation(file_path)
            titles = []
            
            for slide in prs.slides:
                title = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.is_placeholder:
                        placeholder = shape.placeholder_format
                        if placeholder.type == 1:  # Title placeholder
                            title = shape.text.strip()
                            break
                
                titles.append(title if title else f"[Untitled Slide]")
            
            return titles
        
        except Exception as e:
            self.logger.error(f"Error extracting slide titles: {e}")
            return []
